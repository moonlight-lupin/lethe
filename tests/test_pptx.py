"""PowerPoint round-trip: build a deck (title, body, grouped shape, table,
speaker notes), redact it, and verify the output is a working .pptx with no
name leaks and tokens in place."""
import io
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pptx import Presentation
from pptx.util import Inches

from lethe import (Entity, assign_tokens, build_replacer, build_restorer,
                   detect, extract_text, redact_document)

ents = [Entity("John Smith", "PERSON", ["Smith"]),
        Entity("Acme Capital Partners", "COUNTERPARTY", ["Acme"]),
        Entity("Garuda Ventures", "COUNTERPARTY", [])]

# ---- build a deck ------------------------------------------------------------
prs = Presentation()
s1 = prs.slides.add_slide(prs.slide_layouts[1])          # title + content
s1.shapes.title.text = "Project update — Acme Capital Partners"
s1.placeholders[1].text = "Mandate led by John Smith.\nCounter-signed by Garuda Ventures."
s1.notes_slide.notes_text_frame.text = "Remind John Smith about the side letter."

s2 = prs.slides.add_slide(prs.slide_layouts[5])          # blank-ish + table
tbl = s2.shapes.add_table(3, 2, Inches(1), Inches(1), Inches(6), Inches(2)).table
tbl.cell(0, 0).text = "Counterparty"
tbl.cell(0, 1).text = "Lead"
tbl.cell(1, 0).text = "Acme Capital Partners"
tbl.cell(1, 1).text = "John Smith"
tbl.cell(2, 0).text = "Garuda Ventures"
tbl.cell(2, 1).text = "Smith"

buf = io.BytesIO(); prs.save(buf); pptx_bytes = buf.getvalue()

# ---- extract → detect → redact ----------------------------------------------
text = extract_text(pptx_bytes, "pptx")
assert "Slide 1" in text and "Slide 2" in text, "slide markers missing"
assert "Remind John Smith" in text, "speaker notes not extracted"

items = assign_tokens(detect(text, ents))
repl, t2r = build_replacer(items)
out, ext, hits = redact_document(pptx_bytes, "pptx", repl)
print("PPTX ext:", ext, "hits:", hits)
assert ext == ".pptx" and hits > 0

# ---- verify the output deck ---------------------------------------------------
prs2 = Presentation(io.BytesIO(out))                     # loads = valid pptx
out_text = extract_text(out, "pptx")
leaks = [n for n in ("John Smith", "Acme Capital Partners", "Garuda Ventures", "Smith")
         if n in out_text]
print("  leaks:", leaks or "none",
      "| has tokens:", "[PERSON_001]" in out_text and "[COUNTERPARTY_001]" in out_text)
assert not leaks, f"names leaked: {leaks}"
assert "[PERSON_001]" in out_text and "[COUNTERPARTY_001]" in out_text

# table survived as a table (slide 2 still has a graphic frame with a table)
tables = [shp for shp in prs2.slides[1].shapes if getattr(shp, "has_table", False)]
assert tables, "table lost in redaction"
print("  table preserved:", bool(tables))

# round-trip restore on extracted text
restored, n = build_restorer(t2r)(out_text)
print("  restore hits:", n, "| John Smith back:", "John Smith" in restored)
assert "John Smith" in restored and "Acme Capital Partners" in restored
print("PPTX OK")
