"""
Document read / write.

extract_text(...)  -> plain text for the detection + review step.
redact_document(...) -> writes a de-identified file IN THE SAME FORMAT where
                        possible (.docx, .xlsx, .pptx), preserving structure.
                        PDFs and emails (.eml/.msg/.html) are rebuilt as a
                        de-identified .docx because they can't be safely edited in
                        place -- and a clean text/docx is what you feed an AI anyway.
"""
from __future__ import annotations

import contextlib
import email
import functools
import io
import os
import sys
import urllib.request
import xml.etree.ElementTree as ET
import zipfile
from email import policy as email_policy
from html.parser import HTMLParser

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

from . import DATA_DIR


@contextlib.contextmanager
def _suppress_native_stderr():
    """Mute OS-level stderr around a call. Tesseract/PDFium emit warnings ("Image
    too small to scale", vertical-script fallbacks, …) straight to file
    descriptor 2, bypassing Python — this keeps the app's console window clean.
    Python exceptions still propagate normally."""
    try:
        fd = sys.stderr.fileno()
    except (AttributeError, ValueError, io.UnsupportedOperation):
        yield  # no real stderr (e.g. captured) — nothing to mute
        return
    saved = os.dup(fd)
    devnull = os.open(os.devnull, os.O_WRONLY)
    try:
        os.dup2(devnull, fd)
        yield
    finally:
        os.dup2(saved, fd)
        os.close(devnull)
        os.close(saved)

# SpreadsheetML namespace — cell text lives in <si>/<is> elements under this ns.
_SS_NS = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"
_XML_SPACE = "{http://www.w3.org/XML/1998/namespace}space"

try:
    import pdfplumber  # better reading order + table extraction
    _HAS_PDFPLUMBER = True
except Exception:  # pragma: no cover - optional, falls back to pypdf
    _HAS_PDFPLUMBER = False

try:
    import liteparse as _liteparse  # fully-local OCR (PDFium + bundled Tesseract)
    _HAS_LITEPARSE = True
except Exception:  # pragma: no cover - optional; image pages are warned instead
    _HAS_LITEPARSE = False


def ocr_available() -> bool:
    """True when the local OCR engine (liteparse) is installed — scanned /
    image-based PDF pages are then read instead of merely warned about."""
    return _HAS_LITEPARSE


# ---- OCR language data (offline Tesseract models) ----------------------------
# The OCR engine reads its language models (`<code>.traineddata`) from a folder
# under DATA_DIR. English is bundled by the Windows build so OCR works fully
# offline; extra languages are downloaded on demand from the Settings tab
# (alongside that language's name-detection model) into this same folder. The
# engine never reaches the internet at OCR time — it only reads these local files.
# (A pip/pipx install ships no bundled model; liteparse then fetches English from
# tessdata_best on first use, after which it too is offline.)
_TESSDATA_BEST = "https://github.com/tesseract-ocr/tessdata_best/raw/main/{code}.traineddata"


def tessdata_dir() -> str:
    """Folder holding the offline OCR models — `<DATA_DIR>/tessdata`, alongside
    the user's other data (entities.json, vault/). Created on first use."""
    d = os.path.join(DATA_DIR, "tessdata")
    os.makedirs(d, exist_ok=True)
    return d


def installed_ocr_languages() -> list[str]:
    """Primary Tesseract codes to run an OCR pass for (e.g. ['eng', 'jpn']),
    English first. Excludes `*_vert` helper models, which Tesseract loads
    automatically alongside their base language."""
    try:
        codes = [f[:-len(".traineddata")] for f in os.listdir(tessdata_dir())
                 if f.endswith(".traineddata") and not f.endswith("_vert.traineddata")]
    except OSError:
        codes = []
    return (["eng"] if "eng" in codes else []) + sorted(c for c in codes if c != "eng")


def download_ocr_language(codes: list[str]) -> tuple[bool, str]:
    """Fetch one or more Tesseract `tessdata_best` models into the offline
    folder. Needs internet (an explicit, user-initiated Settings action — the
    same moment the spaCy model is fetched). Returns (ok, log)."""
    d = tessdata_dir()
    logs: list[str] = []
    for code in codes:
        dest = os.path.join(d, f"{code}.traineddata")
        if os.path.exists(dest):
            logs.append(f"{code}: already present")
            continue
        tmp = dest + ".part"
        try:
            urllib.request.urlretrieve(_TESSDATA_BEST.format(code=code), tmp)
            os.replace(tmp, dest)
            logs.append(f"{code}: downloaded")
        except Exception as exc:  # noqa: BLE001
            try:
                os.path.exists(tmp) and os.remove(tmp)
            except OSError:
                pass
            return False, "\n".join(logs + [f"{code}: download failed — {exc}"])
    return True, "\n".join(logs)


def remove_ocr_language(codes: list[str]) -> tuple[bool, str]:
    """Delete added OCR models. The bundled English baseline is kept."""
    d = tessdata_dir()
    for code in codes:
        if code == "eng":
            continue
        p = os.path.join(d, f"{code}.traineddata")
        try:
            if os.path.exists(p):
                os.remove(p)
        except OSError as exc:
            return False, f"{code}: remove failed — {exc}"
    return True, "removed"

# A page that carries a raster image yet yields fewer than this many words is
# treated as image-based (a scan, or a figure/chart with baked-in text) — names
# rendered as pixels can't be extracted, so they can't be detected or redacted.
# Gating on an actual image avoids false-flagging legitimately short text pages.
_IMAGE_PAGE_WORD_THRESHOLD = 20


def file_kind(filename: str) -> str:
    ext = os.path.splitext(filename)[1].lower()
    return {".docx": "docx", ".pdf": "pdf", ".xlsx": "xlsx", ".pptx": "pptx",
            ".txt": "txt", ".eml": "eml", ".msg": "msg",
            ".html": "html", ".htm": "html"}.get(ext, "")


# ---- PowerPoint helpers ------------------------------------------------------
def _iter_pptx_shapes(shapes):
    """All shapes on a slide/master/layout, descending into grouped shapes."""
    for shp in shapes:
        if getattr(shp, "shape_type", None) == 6:  # MSO_SHAPE_TYPE.GROUP
            yield from _iter_pptx_shapes(shp.shapes)
        else:
            yield shp


def _iter_pptx_text_frames(prs):
    """Every text frame in the deck: slide shapes (incl. groups), table cells,
    speaker notes, and slide-master/layout shapes (headers, footers, fixed text)."""
    def from_shapes(shapes):
        for shp in _iter_pptx_shapes(shapes):
            if getattr(shp, "has_text_frame", False):
                yield shp.text_frame
            if getattr(shp, "has_table", False):
                for row in shp.table.rows:
                    for cell in row.cells:
                        yield cell.text_frame

    for slide in prs.slides:
        yield from from_shapes(slide.shapes)
        if slide.has_notes_slide:
            yield slide.notes_slide.notes_text_frame
    for master in prs.slide_masters:
        yield from from_shapes(master.shapes)
        for layout in master.slide_layouts:
            yield from from_shapes(layout.shapes)


# ---- PDF page model ---------------------------------------------------------
def _merge_ocr(data: bytes, pages: list[dict]) -> None:
    """OCR the image-based pages in place (when liteparse is installed).

    Pages that carry a raster image but almost no extractable text get their
    text recovered with the fully-local OCR engine (PDFium + bundled Tesseract,
    selective: only the flagged pages are processed) and merged into the page
    narrative, marked `ocr: True` so the UI can ask the user to review them.
    Pages that still come back empty keep the hard "can't read this" warning.
    """
    flagged = [pg for pg in pages
               if pg["n_images"] > 0 and pg["n_words"] < _IMAGE_PAGE_WORD_THRESHOLD]
    if not (flagged and _HAS_LITEPARSE):
        return
    targets = ",".join(str(pg["n"]) for pg in flagged)
    td = tessdata_dir()
    # Models are read from the offline folder, so the engine never reaches the
    # internet at OCR time. This engine has no multi-language ("eng+jpn") mode,
    # so when the user has added languages we OCR each flagged page once per
    # installed language and keep the highest-confidence reading. English-only
    # (the default) is a single pass. Native Tesseract/PDFium warnings are muted.
    langs = installed_ocr_languages() or ["eng"]
    best: dict[int, tuple] = {}  # original page number -> (score, text)
    for lang in langs:
        try:
            with _suppress_native_stderr():
                res = _liteparse.LiteParse(ocr_enabled=True, target_pages=targets, quiet=True,
                                           ocr_language=lang, tessdata_path=td).parse(data)
        except Exception:
            continue  # OCR is best-effort; a failed language is simply skipped
        for pg in flagged:
            # get_page() is keyed by ORIGINAL page number (None outside targets).
            try:
                p = res.get_page(pg["n"])
            except Exception:
                p = None
            text = (p.text or "").strip() if p is not None else ""
            if not text:
                continue
            items = getattr(p, "text_items", None) or []
            confs = [c for it in items if (c := getattr(it, "confidence", None)) is not None]
            score = (sum(confs) / len(confs) if confs else 0.0, len(text))
            if pg["n"] not in best or score > best[pg["n"]][0]:
                best[pg["n"]] = (score, text)
    for pg in flagged:
        if pg["n"] in best:
            text = best[pg["n"]][1]
            pg["narrative"] = f"{pg['narrative']}\n{text}".strip() if pg["narrative"] else text
            pg["n_words"] = len(pg["narrative"].split()) + sum(
                len(c.split()) for tb in pg["tables"] for row in tb for c in row)
            pg["ocr"] = True


@functools.lru_cache(maxsize=4)
def _read_pdf(data: bytes) -> list[dict]:
    """Parse a PDF into per-page content, preserving page boundaries.

    Each page -> {n, narrative, tables, n_words, n_images, ocr} where
    `narrative` is the flowing text *outside* any detected table and `tables`
    is a list of grids (list of rows of cell strings). Uses pdfplumber when
    available for reading order + table detection (flat pypdf text otherwise),
    then OCRs image-based pages when liteparse is installed.

    Cached per document (parsing + OCR are called from extract/warn/redact);
    callers must treat the result as read-only.
    """
    pages = _read_pdf_impl(data)
    _merge_ocr(data, pages)
    return pages


def _read_pdf_impl(data: bytes) -> list[dict]:
    if _HAS_PDFPLUMBER:
        pages: list[dict] = []
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                try:
                    found = page.find_tables()
                except Exception:
                    found = []
                bboxes = [t.bbox for t in found]

                def _outside(obj, _bboxes=bboxes):
                    cx = (obj["x0"] + obj["x1"]) / 2
                    cy = (obj["top"] + obj["bottom"]) / 2
                    return not any(x0 <= cx <= x1 and top <= cy <= bottom
                                   for (x0, top, x1, bottom) in _bboxes)

                try:
                    narrative = (page.filter(_outside).extract_text() if bboxes
                                 else page.extract_text()) or ""
                except Exception:
                    narrative = page.extract_text() or ""

                tables: list[list[list[str]]] = []
                for t in found:
                    try:
                        rows = t.extract()
                    except Exception:
                        continue
                    norm = [[(c or "").strip() for c in row] for row in rows if any(row)]
                    if norm:
                        tables.append(norm)

                n_words = len(narrative.split()) + sum(
                    len(c.split()) for tb in tables for row in tb for c in row)
                pages.append({"n": i, "narrative": narrative, "tables": tables,
                              "n_words": n_words, "n_images": len(page.images or []),
                              "ocr": False})
        return pages

    # Fallback: pypdf flat text, page by page (no tables, page anchors kept).
    reader = PdfReader(io.BytesIO(data))
    out = []
    for i, page in enumerate(reader.pages, 1):
        txt = page.extract_text() or ""
        out.append({"n": i, "narrative": txt, "tables": [],
                    "n_words": len(txt.split()), "n_images": 0, "ocr": False})
    return out


def pdf_warnings(data: bytes) -> list[dict]:
    """Pages needing user attention, two kinds:
    - ocr=True  — the page was image-based and its text was recovered with the
      local OCR engine; it IS detected/redacted, but OCR isn't perfect → review.
    - ocr=False — the page is image-based and could NOT be read (no OCR engine
      installed, or OCR found nothing); names there are NOT redacted."""
    out = []
    for pg in _read_pdf(data):
        if pg.get("ocr"):
            out.append({"page": pg["n"], "words": pg["n_words"],
                        "images": pg["n_images"], "ocr": True})
        elif pg["n_images"] > 0 and pg["n_words"] < _IMAGE_PAGE_WORD_THRESHOLD:
            out.append({"page": pg["n"], "words": pg["n_words"],
                        "images": pg["n_images"], "ocr": False})
    return out


# ---- extraction (for detection) ---------------------------------------------
def extract_text(data: bytes, kind: str) -> str:
    if kind == "txt":
        return data.decode("utf-8", errors="replace")
    if kind == "docx":
        doc = Document(io.BytesIO(data))
        parts = [p.text for p in doc.paragraphs]
        for table in doc.tables:
            for row in table.rows:
                parts.extend(cell.text for cell in row.cells)
        for section in doc.sections:
            for hdrftr in (section.header, section.footer):
                parts.extend(p.text for p in hdrftr.paragraphs)
        return "\n".join(parts)
    if kind == "xlsx":
        wb = load_workbook(io.BytesIO(data), data_only=True, read_only=True)
        parts = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                parts.extend(str(c) for c in row if c is not None)
        return "\n".join(parts)
    if kind == "pdf":
        parts = []
        for pg in _read_pdf(data):
            parts.append(f"Page {pg['n']}")
            if pg["narrative"]:
                parts.append(pg["narrative"])
            for tb in pg["tables"]:
                for row in tb:
                    parts.append("\t".join(row))
        return "\n".join(parts)
    if kind == "pptx":
        prs = Presentation(io.BytesIO(data))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"Slide {i}")
            for shp in _iter_pptx_shapes(slide.shapes):
                if getattr(shp, "has_text_frame", False) and shp.text_frame.text.strip():
                    parts.append(shp.text_frame.text)
                if getattr(shp, "has_table", False):
                    for row in shp.table.rows:
                        parts.append("\t".join(cell.text for cell in row.cells))
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                parts.append("Notes: " + slide.notes_slide.notes_text_frame.text)
        master_bits = []
        for master in prs.slide_masters:
            for shapes in [master.shapes] + [lo.shapes for lo in master.slide_layouts]:
                for shp in _iter_pptx_shapes(shapes):
                    if getattr(shp, "has_text_frame", False) and shp.text_frame.text.strip():
                        master_bits.append(shp.text_frame.text)
        if master_bits:
            parts.append("Master / layout text")
            parts.extend(master_bits)
        return "\n".join(parts)
    if kind in ("eml", "msg", "html"):
        parsed = _parse_email(data, kind)
        parts = [f"{label}: {val}" for label, val in parsed["headers"]]
        if parsed["body"]:
            parts.append(parsed["body"])
        return "\n".join(parts)
    raise ValueError(f"Unsupported kind: {kind}")


def read_xlsx_grid(data: bytes, max_rows: int = 200, max_cols: int = 40):
    """Read a workbook as a grid for the live preview: a list of
    (sheet_name, rows, truncated) where rows is a list of lists of cell strings.
    Capped per sheet so a huge workbook stays responsive."""
    wb = load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    sheets = []
    try:
        for ws in wb.worksheets:
            rows, truncated = [], False
            for ri, row in enumerate(ws.iter_rows(values_only=True)):
                if ri >= max_rows:
                    truncated = True
                    break
                if len(row) > max_cols:
                    truncated = True
                rows.append(["" if c is None else str(c) for c in row[:max_cols]])
            while rows and not any(c.strip() for c in rows[-1]):
                rows.pop()                       # trim trailing empty rows
            sheets.append((ws.title, rows, truncated))
    finally:
        wb.close()
    return sheets


# ---- redacted output --------------------------------------------------------
def _redact_paragraph(paragraph, replace_fn) -> int:
    """Replace text in a paragraph. If anything changes we flatten the
    paragraph's runs into one (formatting within a changed line is lost, but the
    redaction is guaranteed correct -- the right trade-off for a privacy gate)."""
    original = paragraph.text
    if not original:
        return 0
    new, hits = replace_fn(original)
    if hits and new != original:
        for run in list(paragraph.runs):
            run.text = ""
        if paragraph.runs:
            paragraph.runs[0].text = new
        else:
            paragraph.add_run(new)
    return hits


def _redact_docx(data: bytes, replace_fn) -> tuple[bytes, int]:
    doc = Document(io.BytesIO(data))
    hits = 0
    for p in doc.paragraphs:
        hits += _redact_paragraph(p, replace_fn)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    hits += _redact_paragraph(p, replace_fn)
    for section in doc.sections:
        for hdrftr in (section.header, section.footer):
            for p in hdrftr.paragraphs:
                hits += _redact_paragraph(p, replace_fn)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue(), hits


def _redact_string_table(xml_bytes: bytes, replace_fn, container: str):
    """Redact the text inside every <si> (shared strings) or <is> (inline
    strings) element of an OOXML part, leaving the rest of the XML structure
    intact. A string whose text changes is collapsed to a single <t> (any
    in-cell rich-text formatting on that one string is lost — the same trade-off
    we make in Word). Returns (new_bytes, hits) or (None, 0) if nothing changed
    or the part can't be parsed."""
    try:
        ET.register_namespace("", _SS_NS)
        root = ET.fromstring(xml_bytes)
    except ET.ParseError:
        return None, 0
    t_tag = f"{{{_SS_NS}}}t"
    cont_tag = f"{{{_SS_NS}}}{container}"
    hits = 0
    for cont in root.iter(cont_tag):
        t_elems = cont.findall(f".//{t_tag}")
        full = "".join(t.text or "" for t in t_elems)
        if not full:
            continue
        new, n = replace_fn(full)
        if n and new != full:
            hits += n
            for child in list(cont):          # drop existing runs/text
                cont.remove(child)
            t = ET.SubElement(cont, t_tag)     # replace with one plain <t>
            t.set(_XML_SPACE, "preserve")
            t.text = new
    if not hits:
        return None, 0
    body = ET.tostring(root, encoding="unicode")
    out = '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\r\n' + body
    return out.encode("utf-8"), hits


def _redact_xlsx(data: bytes, replace_fn) -> tuple[bytes, int]:
    """Redact an .xlsx by surgically editing only the string parts of the OOXML
    package and copying every other part byte-for-byte. Unlike a full openpyxl
    load/save, this preserves charts, images, pivot tables and formatting — so
    Excel doesn't flag the file for repair."""
    hits = 0
    out_buf = io.BytesIO()
    with zipfile.ZipFile(io.BytesIO(data), "r") as src, \
            zipfile.ZipFile(out_buf, "w", zipfile.ZIP_DEFLATED) as dst:
        for item in src.infolist():
            content = src.read(item.filename)
            name = item.filename
            if name == "xl/sharedStrings.xml":
                new, n = _redact_string_table(content, replace_fn, "si")
                if new is not None:
                    content, hits = new, hits + n
            elif name.startswith("xl/worksheets/") and name.endswith(".xml"):
                new, n = _redact_string_table(content, replace_fn, "is")
                if new is not None:
                    content, hits = new, hits + n
            dst.writestr(item, content)          # ZipInfo preserves metadata
    return out_buf.getvalue(), hits


def _redact_pptx(data: bytes, replace_fn) -> tuple[bytes, int]:
    """Redact a PowerPoint deck in place — slide text (incl. grouped shapes),
    tables, speaker notes and master/layout text — and return a working .pptx.
    As in Word, a paragraph whose text changes is flattened to a single run
    (in-line mixed formatting within that one paragraph is lost; the redaction
    is guaranteed correct)."""
    prs = Presentation(io.BytesIO(data))
    hits = 0
    for tf in _iter_pptx_text_frames(prs):
        for p in tf.paragraphs:
            original = "".join(run.text for run in p.runs)
            if not original:
                continue
            new, n = replace_fn(original)
            if n and new != original:
                hits += n
                for run in p.runs:
                    run.text = ""
                p.runs[0].text = new
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue(), hits


def _text_to_docx(text: str) -> bytes:
    doc = Document()
    for line in text.split("\n"):
        doc.add_paragraph(line)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _add_notice_header(doc) -> None:
    """Prepend an agent/reader-facing notice: what this document is, and — most
    importantly — how to cite it (by the ORIGINAL PDF page, via the `Page N`
    headings, not this Word file's own pagination)."""
    doc.add_heading("Document notice — converted from PDF & de-identified", level=1)
    doc.add_paragraph(
        "This document was automatically converted from a PDF and de-identified by "
        "Lethe before it reached you. If you are an AI agent or tool processing this "
        "file, read these notes before quoting or analysing it:")
    for note in (
        "CITE BY ORIGINAL PDF PAGE. The body below is divided by “Page N” headings, "
        "where N is the page number in the source PDF. When you quote, reference or "
        "footnote anything, use that source page number — NOT this Word document's own "
        "pagination, which does not match the original PDF.",
        "Redacted identities appear as placeholder tokens such as [PERSON_001] or "
        "[COUNTERPARTY_001]. Treat each token as one opaque identifier and reproduce it "
        "verbatim — do not alter, expand, translate or guess the real name behind it.",
        "Some content may be missing: text inside images, charts, logos or scanned pages "
        "is not captured (there is no OCR), so any names or figures there are neither "
        "shown nor redacted here.",
        "Tables were reconstructed from the PDF; fine formatting and exact positioning "
        "may differ from the original.",
    ):
        doc.add_paragraph(note, style="List Bullet")
    doc.add_paragraph("—" * 24)


def _pdf_to_docx(pages: list[dict], replace_fn) -> tuple[bytes, int]:
    """Rebuild a de-identified .docx from parsed PDF pages.

    Starts with an agent-facing notice header (how to cite, what the tokens mean,
    what's missing), then preserves page boundaries (a `Page N` heading + page
    break per source page) and renders detected tables as real Word tables. The
    `Page N` anchors let downstream tools/agents quote against the original PDF
    pages."""
    doc = Document()
    _add_notice_header(doc)

    hits = 0
    for idx, pg in enumerate(pages):
        if idx > 0:
            doc.add_page_break()
        doc.add_heading(f"Page {pg['n']}", level=2)

        if pg["narrative"]:
            new, n = replace_fn(pg["narrative"])
            hits += n
            for line in new.split("\n"):
                doc.add_paragraph(line)

        for tb in pg["tables"]:
            ncols = max((len(row) for row in tb), default=0)
            if not ncols:
                continue
            table = doc.add_table(rows=0, cols=ncols)
            try:
                table.style = "Table Grid"
            except Exception:
                pass
            for row in tb:
                cells = table.add_row().cells
                for ci in range(ncols):
                    val = row[ci] if ci < len(row) else ""
                    if val:
                        new, n = replace_fn(val)
                        hits += n
                        cells[ci].text = new
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue(), hits


# ---- Email (.eml / .msg) & HTML ---------------------------------------------
# Emails and saved-as-HTML messages can't be safely round-tripped, so — like
# PDFs — they come in and a de-identified Word file goes out. We pull the header
# block (From / To / Cc / Subject) and the message body; the headers carry the
# very names and addresses Lethe targets, so they are redacted with the body.
_HTML_BLOCK = {"p", "div", "br", "li", "tr", "table", "h1", "h2", "h3", "h4",
               "h5", "h6", "blockquote", "ul", "ol", "section", "article",
               "header", "footer", "hr"}
_HTML_SKIP = {"script", "style", "head", "title"}


class _HTMLTextExtractor(HTMLParser):
    """Minimal, dependency-free HTML → text: drops tags, turns block elements
    into line breaks, skips <script>/<style>. Good enough for de-identifying an
    email body (we want the words, not a faithful re-render)."""

    def __init__(self):
        super().__init__(convert_charrefs=True)
        self._out: list[str] = []
        self._skip = 0

    def handle_starttag(self, tag, attrs):
        if tag in _HTML_SKIP:
            self._skip += 1
        elif tag in _HTML_BLOCK:
            self._out.append("\n")

    def handle_endtag(self, tag):
        if tag in _HTML_SKIP and self._skip:
            self._skip -= 1
        elif tag in _HTML_BLOCK:
            self._out.append("\n")

    def handle_data(self, data):
        if not self._skip:
            self._out.append(data)

    def text(self) -> str:
        return "".join(self._out)


def _html_to_text(html: str) -> str:
    p = _HTMLTextExtractor()
    try:
        p.feed(html)
    except Exception:  # noqa: BLE001 — malformed HTML shouldn't crash de-id
        pass
    lines = [ln.strip() for ln in p.text().splitlines()]
    out: list[str] = []
    blank = False
    for ln in lines:  # collapse runs of blank lines
        if ln:
            out.append(ln)
            blank = False
        elif not blank:
            out.append("")
            blank = True
    return "\n".join(out).strip()


_EMAIL_HEADERS = [("From", "from"), ("To", "to"), ("Cc", "cc"),
                  ("Date", "date"), ("Subject", "subject")]


def _parse_eml(data: bytes) -> dict:
    msg = email.message_from_bytes(data, policy=email_policy.default)
    headers = [(label, str(msg[key])) for label, key in _EMAIL_HEADERS if msg[key]]
    body = ""
    try:
        part = msg.get_body(preferencelist=("plain", "html")) if msg.is_multipart() else msg
        if part is not None:
            content = part.get_content()
            if isinstance(content, bytes):
                content = content.decode("utf-8", errors="replace")
            body = _html_to_text(content) if part.get_content_subtype() == "html" else content
    except Exception:  # noqa: BLE001 — odd MIME shouldn't crash de-id
        body = ""
    return {"headers": headers, "body": body}


def _parse_msg(data: bytes) -> dict:
    try:
        import extract_msg
    except ImportError as exc:  # the optional [email] extra isn't installed
        raise ValueError(
            "Reading Outlook .msg files needs the optional 'email' extra — "
            "install it with  pip install \"lethe[email]\"  (or use the Windows "
            "installer / portable bundle, which include it). .eml and .html files "
            "work without it.") from exc
    m = extract_msg.openMsg(data)
    fields = [("From", m.sender), ("To", m.to), ("Cc", m.cc),
              ("Date", m.date), ("Subject", m.subject)]
    headers = [(label, str(v)) for label, v in fields if v]
    body = m.body or ""
    if not body and getattr(m, "htmlBody", None):
        html = m.htmlBody
        if isinstance(html, bytes):
            html = html.decode("utf-8", errors="replace")
        body = _html_to_text(html)
    return {"headers": headers, "body": str(body)}


def _parse_email(data: bytes, kind: str) -> dict:
    if kind == "eml":
        return _parse_eml(data)
    if kind == "msg":
        return _parse_msg(data)
    # html / htm — a saved web page or exported message; no separate header block
    return {"headers": [], "body": _html_to_text(data.decode("utf-8", errors="replace"))}


def _add_email_notice(doc) -> None:
    doc.add_heading("Document notice — converted from an email & de-identified", level=1)
    doc.add_paragraph(
        "This message was automatically converted from an email / HTML file and "
        "de-identified by Lethe before it reached you. If you are an AI agent or tool "
        "processing this file, read these notes first:")
    for note in (
        "Redacted identities appear as placeholder tokens such as [PERSON_001] or "
        "[COUNTERPARTY_001]. Treat each token as one opaque identifier and reproduce it "
        "verbatim — do not alter, expand, translate or guess the real name behind it.",
        "The header block (From / To / Cc / Subject) has been de-identified along with the "
        "message body.",
        "Attachments and inline images are NOT included or de-identified — only the message "
        "text is. Any names inside an attachment or image are neither shown nor redacted here.",
        "Formatting is simplified: the original HTML layout, fonts and tables are rendered as "
        "plain text.",
    ):
        doc.add_paragraph(note, style="List Bullet")
    doc.add_paragraph("—" * 24)


def _email_to_docx(parsed: dict, replace_fn) -> tuple[bytes, int]:
    """Rebuild a de-identified .docx from a parsed email/HTML message: an
    agent-facing notice, the labelled From/To/Cc/Date/Subject block, then body."""
    doc = Document()
    _add_email_notice(doc)
    hits = 0
    for label, val in parsed["headers"]:
        new, n = replace_fn(val)
        hits += n
        p = doc.add_paragraph()
        p.add_run(f"{label}: ").bold = True
        p.add_run(new)
    if parsed["headers"]:
        doc.add_paragraph("—" * 24)
    for line in (parsed["body"] or "").split("\n"):
        new, n = replace_fn(line)
        hits += n
        doc.add_paragraph(new)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue(), hits


def redact_document(data: bytes, kind: str, replace_fn) -> tuple[bytes, str, int]:
    """Return (output_bytes, output_extension, hit_count)."""
    if kind == "docx":
        out, hits = _redact_docx(data, replace_fn)
        return out, ".docx", hits
    if kind == "xlsx":
        out, hits = _redact_xlsx(data, replace_fn)
        return out, ".xlsx", hits
    if kind == "txt":
        new, hits = replace_fn(data.decode("utf-8", errors="replace"))
        return new.encode("utf-8"), ".txt", hits
    if kind == "pptx":
        out, hits = _redact_pptx(data, replace_fn)
        return out, ".pptx", hits
    if kind == "pdf":
        out, hits = _pdf_to_docx(_read_pdf(data), replace_fn)
        return out, ".docx", hits
    if kind in ("eml", "msg", "html"):
        out, hits = _email_to_docx(_parse_email(data, kind), replace_fn)
        return out, ".docx", hits
    raise ValueError(f"Unsupported kind: {kind}")
